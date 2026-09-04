#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""答辩 PPT 半自动生成器（reveal.js HTML 或 python-pptx）。

用法:
    # 生成 reveal.js HTML（浏览器打开即可放映，无需装软件）
    python gen_defense.py --format revealjs --out report/defense.html

    # 生成 PPTX（需 pip install python-pptx）
    python gen_defense.py --format pptx --out report/defense.pptx

数据源:
    - out/results.json        关键数字
    - report/main.tex/.docx   论文结构（可选，提取章节标题）
    - out/runs.csv            实验追踪记录（可选，生成模型演进图）

输出结构（10-12 页）:
    1. 封面：题目、队名、成员、日期
    2. 问题背景与重述（1 页）
    3. 核心建模思路总览（1 页，流程图占位）
    4. 问题一：模型 + 关键结果 + 图表
    5. 问题二：模型 + 关键结果 + 图表
    6. 问题三：模型 + 关键结果 + 图表
    7. 模型检验与灵敏度（灵敏度表/图）
    8. 创新亮点（2-3 条，有证据）
    9. 局限与改进（诚实、机制化）
    10. 落点与建议（可执行、面向决策者）
    11. 团队分工与复盘（可选）
    12. Q&A 备选页（预判评委问点）
"""
import argparse
import json
import os
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional

# ==================== 数据加载 ====================

def load_results(results_path: str = "out/results.json") -> Dict[str, Any]:
    p = Path(results_path)
    if not p.exists():
        return {}
    with open(p, encoding="utf-8") as f:
        return json.load(f)


def load_runs_csv(runs_path: str = "out/runs.csv") -> List[Dict[str, Any]]:
    p = Path(runs_path)
    if not p.exists():
        return []
    import csv
    rows = []
    with open(p, encoding="utf-8") as f:
        reader = csv.DictReader(f)
        for row in reader:
            try:
                row["params"] = json.loads(row["params_json"])
                row["metrics"] = json.loads(row["metrics_json"])
            except Exception:
                row["params"] = {}
                row["metrics"] = {}
            rows.append(row)
    return rows


def extract_paper_sections(tex_path: str = "report/main.tex") -> List[str]:
    """从 LaTeX 提取 section 标题作为大纲参考"""
    p = Path(tex_path)
    if not p.exists():
        return []
    import re
    text = p.read_text(encoding="utf-8", errors="ignore")
    sections = re.findall(r"\\section\{(.*?)\}", text)
    subsections = re.findall(r"\\subsection\{(.*?)\}", text)
    return sections + subsections


# ==================== 内容模板 ====================

def build_slides_data(results: Dict, runs: List[Dict], paper_sections: List[str]) -> List[Dict]:
    """构建幻灯片数据结构，每项是一页"""
    slides = []

    # 1. 封面
    slides.append({
        "type": "cover",
        "title": "基于<模型名>的<对象><任务>",
        "subtitle": "数学建模竞赛答辩汇报",
        "meta": ["队名：<队名>", "成员：<成员1、成员2、成员3>", "日期：<答辩日期>"],
    })

    # 2. 问题背景
    slides.append({
        "type": "content",
        "title": "问题背景与重述",
        "bullets": [
            "背景：<一句话背景，含关键统计/引用>",
            "核心任务：",
            "  • 问题一：<任务描述>",
            "  • 问题二：<任务描述>",
            "  • 问题三：<任务描述>",
            "关键约束/数据特点：<约束条件、数据规模、特殊性质>",
        ],
    })

    # 3. 总览
    slides.append({
        "type": "content",
        "title": "核心建模思路总览",
        "bullets": [
            "整体流程：数据清洗 → 假设建立 → 分问建模 → 求解验证 → 论文撰写",
            "建模框架图：<此处放 build_modeling_flowchart.png / 手绘拍照>",
            "关键创新点前瞻：",
            "  ① 多模型对照择优 / 机理+数据耦合 / 问题重构 …",
            "  ② 一个扎实的理论性质（唯一解/凸性/误差界）",
            "  ③ 检验做成证据链（蒙特卡洛/最坏情形/基准对比）",
        ],
    })

    # 4-6. 每问详情
    for q_num in [1, 2, 3]:
        q_key = f"q{q_num}"
        model_name = results.get(f"{q_key}_model", f"<问题{q_num}模型名>")
        opt_val = results.get(f"{q_key}_optimal_value", results.get(f"{q_key}_opt", "<数值>"))
        error = results.get(f"{q_key}_error_pct", results.get(f"{q_key}_mape", "<误差>"))

        slides.append({
            "type": "qdetail",
            "title": f"问题{q_num}：{model_name}",
            "model_name": model_name,
            "bullets": [
                f"模型类型：<评价/预测/优化/机理/统计/图论…>",
                f"核心公式/指标：<关键方程或评分体系>",
                f"求解方法：<算法名/库/关键参数>",
                f"关键结果：最优值 = {opt_val}，误差 = {error}",
                f"结果图：<此处放 out/figures/q{q_num}_result.png>",
            ],
            "figure_hint": f"figures/q{q_num}_result.png",
        })

    # 7. 检验与灵敏度
    slides.append({
        "type": "content",
        "title": "模型检验与灵敏度分析",
        "bullets": [
            "量纲/量级/边界校验：已通过 sanity_check.py 全自动校验 ✓",
            "灵敏度分析：关键参数 ±10%/±20% 摄动",
            "  • 最敏感参数：<参数名>，结果变化 <X>% → <结论稳健/需关注>",
            "  • 灵敏度图：<此处放 out/figures/sensitivity.png>",
            "误差分析/交叉验证：<RMSE/MAE/R²/置信区间>",
            "基准对比：<如有> 与 <基准方法> 对比，提升 <X>%",
        ],
    })

    # 8. 创新亮点
    slides.append({
        "type": "content",
        "title": "创新亮点（2–3 条扎实的，有证据支撑）",
        "bullets": [
            "亮点 1：多模型对照择优 —— 同时跑 <模型A/B/C>，用 <指标> 择优，"
            "论文给出对比表与择优理由（见 innovation-playbook.md 套路 #1）",
            "亮点 2：<一处理论改进/性质证明/自定义指标> —— "
            "<具体做了什么 + 为什么有效 + 对结果的量化贡献>（套路 #3/4/7）",
            "亮点 3：检验证据链 —— 蒙特卡洛随机扰动 + 最坏情形 + 基准对比，"
            "画出「参数—结果」证据图，逐个解释敏感参数物理意义（套路 #5）",
        ],
    })

    # 9. 局限
    slides.append({
        "type": "content",
        "title": "局限与改进（诚实、机制化）",
        "bullets": [
            "局限 1：<具体局限，如 假设 X 仅在条件 Y 下成立> —— "
            "机制：<为什么会这样> → 后果：<结果可能偏乐观/保守> → "
            "补救：<未来可尝试的改进方向>",
            "局限 2：<数据规模/质量/计算资源限制> —— "
            "机制 + 后果 + 补救",
            "负结果展示：<如有> 尝试了 <方法> 但 <失败表现>，原因 <机制解释>，"
            "这反而划清了模型适用边界",
        ],
    })

    # 10. 落点
    slides.append({
        "type": "content",
        "title": "落点与建议（面向决策者，可执行）",
        "bullets": [
            "建议 1：对 <决策对象A> 实施 <具体行动>，预期 <量化效果>，"
            "依据：<模型结果/文献/实证>",
            "建议 2：对 <决策对象B> 实施 <具体行动>，预期 <量化效果>…",
            "推广性：该 <模型名/方法论> 可迁移至 <同类场景>，"
            "只需替换 <数据/参数> 即可复现",
        ],
    })

    # 11. 团队分工（可选）
    if runs:
        slides.append({
            "type": "content",
            "title": "团队分工与模型演进",
            "bullets": [
                "分工：建模手（假设/选型/推导）↔ 编程手（清洗/求解/检验）↔ 写作手（摘要/图表/降重）",
                f"实验记录：共 {len(runs)} 次运行，覆盖 {len(set(r['question'] for r in runs))} 个问题、"
                f"{len(set(r['model'] for r in runs))} 个模型",
                "模型演进：<如 q1 从 Baseline → Advanced → Outstanding 三档迭代，关键决策节点>",
                "赛后复盘：已跑 review_survey.py 生成四维复盘草稿",
            ],
        })

    # 12. Q&A 备选
    slides.append({
        "type": "content",
        "title": "Q&A 备选：评委高频问点预判",
        "bullets": [
            "Q1：为什么选这个模型，没考虑 <其他经典模型> 吗？"
            " → A：我们对比了 <模型A/B/C>，<指标> 上 <模型名> 最优，且 <理论理由>…",
            "Q2：假设 <关键假设> 不成立会怎样？"
            " → A：我们做了 <摄动/反事实> 测试，结论 <稳健/会变化 X%>，"
            "适用边界是 <条件>…",
            "Q3：结果的置信度如何？有没有交叉验证？"
            " → A：<交叉验证指标/Bootstrap 区间/蒙特卡洛分布>，"
            "置信区间 <[L, U]>，误差 <X>%",
            "Q4：创新点到底在哪？不是套用现成模型吗？"
            " → A：创新在于 <具体改进/融合/指标/证据链>，并非简单套用，"
            "对结果贡献 <量化说明>…",
            "Q5：代码能不能现场跑一遍？"
            " → A：可以，固定种子、锁版本 requirements.txt、canonical 脚本在 src/export_results.py",
        ],
    })

    return slides


# ==================== 输出格式 ====================

# reveal.js HTML 模板
REVEAL_TEMPLATE = """<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>{title}</title>
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@4.5.0/dist/reveal.css">
    <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@4.5.0/dist/theme/white.css">
    <style>
        .reveal {{ font-family: "Microsoft YaHei", "SimHei", sans-serif; }}
        .reveal h1 {{ font-size: 2.5em; }}
        .reveal h2 {{ font-size: 1.8em; color: #2c3e50; }}
        .reveal ul {{ text-align: left; line-height: 1.6; }}
        .reveal li {{ margin: 0.5em 0; }}
        .cover {{ text-align: center; }}
        .cover h1 {{ font-size: 3em; margin-bottom: 0.5em; }}
        .cover .meta {{ font-size: 1.2em; color: #555; margin-top: 1em; }}
        .qdetail .model-name {{ color: #2980b9; font-weight: bold; }}
        .figure-hint {{ color: #888; font-size: 0.8em; font-style: italic; }}
        .reveal pre code {{ font-size: 0.7em; }}
        .reveal .small {{ font-size: 0.8em; color: #666; }}
    </style>
</head>
<body>
<div class="reveal">
<div class="slides">
{slides_html}
</div>
</div>
<script src="https://cdn.jsdelivr.net/npm/reveal.js@4.5.0/dist/reveal.js"></script>
<script>
    Reveal.initialize({{
        hash: true,
        slideNumber: true,
        transition: 'slide',
        width: 1280,
        height: 720,
        margin: 0.1,
    }});
</script>
</body>
</html>
"""

def render_reveal(slides: List[Dict], title: str) -> str:
    slides_html = []
    for i, s in enumerate(slides):
        cls = s.get("type", "content")
        html = [f'<section data-markdown data-separator="^---$" class="{cls}">']

        if cls == "cover":
            html.append(f"<h1>{s['title']}</h1>")
            if s.get("subtitle"):
                html.append(f"<h3>{s['subtitle']}</h3>")
            if s.get("meta"):
                html.append('<div class="meta">')
                for m in s["meta"]:
                    html.append(f"<p>{m}</p>")
                html.append("</div>")
        else:
            html.append(f"<h2>{s['title']}</h2>")
            if s.get("model_name"):
                html.append(f'<p class="model-name">模型：{s["model_name"]}</p>')
            if s.get("bullets"):
                html.append("<ul>")
                for b in s["bullets"]:
                    html.append(f"<li>{b}</li>")
                html.append("</ul>")
            if s.get("figure_hint"):
                html.append(f'<p class="figure-hint">📊 建议插入图表：{s["figure_hint"]}</p>')

        html.append("</section>")
        slides_html.append("\n".join(html))

    return REVEAL_TEMPLATE.format(title=title, slides_html="\n".join(slides_html))


# python-pptx 生成（可选依赖）
def render_pptx(slides: List[Dict], output_path: str):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        print("[gen_defense] 需安装 python-pptx: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 颜色
    DARK_BLUE = RGBColor(0x2C, 0x3E, 0x50)
    ACCENT = RGBColor(0x29, 0x80, 0xB9)
    GRAY = RGBColor(0x7F, 0x8C, 0x8D)

    def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_BLUE, alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = alignment
        return txBox

    def add_bullets(slide, left, top, width, height, bullets, font_size=16, color=DARK_BLUE, spacing=Pt(6)):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = b.lstrip("• ").lstrip("- ")
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.space_after = spacing
            p.level = b.count("  ") // 2  # 缩进层级
        return txBox

    for s in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        cls = s.get("type", "content")

        if cls == "cover":
            add_textbox(slide, 1, 1.5, 11, 1.5, s["title"], 44, True, DARK_BLUE, PP_ALIGN.CENTER)
            if s.get("subtitle"):
                add_textbox(slide, 1, 3.0, 11, 0.8, s["subtitle"], 24, False, GRAY, PP_ALIGN.CENTER)
            if s.get("meta"):
                y = 4.0
                for m in s["meta"]:
                    add_textbox(slide, 1, y, 11, 0.5, m, 18, False, GRAY, PP_ALIGN.CENTER)
                    y += 0.5
        else:
            add_textbox(slide, 0.5, 0.3, 12, 0.7, s["title"], 28, True, DARK_BLUE)
            if s.get("model_name"):
                add_textbox(slide, 0.5, 1.0, 12, 0.5, f"模型：{s['model_name']}", 16, True, ACCENT)
            if s.get("bullets"):
                add_bullets(slide, 0.5, 1.3, 12, 5.5, s["bullets"], 14, DARK_BLUE)
            if s.get("figure_hint"):
                add_textbox(slide, 0.5, 6.5, 12, 0.5, f"📊 建议插入图表：{s['figure_hint']}", 12, False, GRAY)

        # 页码
        add_textbox(slide, 12, 7.0, 1, 0.3, f"{slides.index(s)+1}/{len(slides)}", 10, False, GRAY, PP_ALIGN.RIGHT)

    prs.save(output_path)
    print(f"[gen_defense] 已生成 PPTX: {output_path}")


# ==================== CLI ====================

def main():
    ap = argparse.ArgumentParser(description="答辩 PPT 半自动生成器")
    ap.add_argument("--format", choices=["revealjs", "pptx"], default="revealjs",
                    help="输出格式：revealjs(HTML, 无依赖) 或 pptx(需 python-pptx)")
    ap.add_argument("--out", "-o", default=None, help="输出路径，默认 report/defense.html 或 .pptx")
    ap.add_argument("--results", default="out/results.json", help="results.json 路径")
    ap.add_argument("--runs", default="out/runs.csv", help="runs.csv 路径")
    ap.add_argument("--paper", default="report/main.tex", help="论文源文件路径（提取大纲）")
    ap.add_argument("--title", default="数学建模竞赛答辩汇报", help="PPT 标题")
    args = ap.parse_args()

    # 加载数据
    results = load_results(args.results)
    runs = load_runs_csv(args.runs)
    paper_sections = extract_paper_sections(args.paper)

    # 构建幻灯片数据
    slides = build_slides_data(results, runs, paper_sections)

    # 确定输出路径
    if args.out:
        out_path = args.out
    else:
        ext = ".html" if args.format == "revealjs" else ".pptx"
        out_path = f"report/defense{ext}"

    os.makedirs(os.path.dirname(out_path), exist_ok=True)

    # 生成
    if args.format == "revealjs":
        html = render_reveal(slides, args.title)
        with open(out_path, "w", encoding="utf-8") as f:
            f.write(html)
        print(f"[gen_defense] 已生成 reveal.js HTML: {out_path}")
        print(f"       用浏览器打开即可放映（按空格翻页，S 键查看演讲者备注）")
    else:
        render_pptx(slides, out_path)


if __name__ == "__main__":
    raise SystemExit(main())